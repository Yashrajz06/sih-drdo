#!/usr/bin/env python3
"""Build docs/references_slide.pptx -- one slide, three cards, live hyperlinks.

Every paper, dataset and standard name is a clickable link. Card heights are
measured from their content rather than fixed, so no card ends in dead space.

LINK POLICY. Entries verified against this repository -- README.md's attribution
section, third_party/gtcrn/README.md, scripts/esc50_noise.py -- get direct URLs.
Papers whose DOI could not be verified get a Google Scholar search link instead
of a guessed one: a search link always resolves to the right paper, whereas a
wrong DOI 404s in front of judges. Swap those for real DOIs once checked.

    .venv/bin/python scripts/make_references_pptx.py
"""
from pathlib import Path
from urllib.parse import quote_plus

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

SCH = "https://scholar.google.com/scholar?q="
sch = lambda q: SCH + quote_plus(q)

INK, MUTED, FAINT = RGBColor(0x0F, 0x17, 0x2A), RGBColor(0x47, 0x55, 0x69), RGBColor(0x7A, 0x87, 0x98)
THEME = {
    "blue":  dict(dark=RGBColor(0x1E, 0x3A, 0x8A), pill=RGBColor(0xBF, 0xDB, 0xFE),
                  body=RGBColor(0xEF, 0xF5, 0xFF), line=RGBColor(0x9F, 0xC0, 0xEE)),
    "amber": dict(dark=RGBColor(0x7C, 0x4A, 0x03), pill=RGBColor(0xFB, 0xD8, 0x7F),
                  body=RGBColor(0xFF, 0xFB, 0xEF), line=RGBColor(0xE5, 0xB8, 0x5A)),
    "green": dict(dark=RGBColor(0x14, 0x53, 0x2D), pill=RGBColor(0xC7, 0xEF, 0xA6),
                  body=RGBColor(0xF4, 0xFC, 0xEE), line=RGBColor(0x9F, 0xD3, 0x7C)),
}

CARDS = [
    ("blue", "Model and Method", [
        ("e", "GTCRN", "https://ieeexplore.ieee.org/document/10448310",
         "Rong et al., ICASSP 2024", "github.com/Xiaobin-Rong/gtcrn  ·  MIT licence"),
        ("e", "Complex Ratio Masking",
         sch("Complex Ratio Masking for Monaural Speech Separation Williamson Wang"),
         "Williamson, Wang & Wang, IEEE/ACM TASLP 24(3), 2016", "the complex mask our model outputs"),
        ("e", "ERB auditory filters",
         sch("Derivation of auditory filter shapes from notched-noise data Glasberg Moore"),
         "Glasberg & Moore, Hearing Research 47, 1990", "the perceptual sub-band front end"),
        ("h", "LOSS FUNCTIONS IMPLEMENTED FROM", "", "", ""),
        ("e", "Braun & Tashev, 2021",
         sch("A consolidated view of loss functions for supervised deep learning based speech enhancement"),
         "compressed magnitude, asymmetric anti-over-suppression", ""),
        ("e", "Yamamoto, Song & Kim, ICASSP 2020", "https://arxiv.org/abs/1910.11480",
         "multi-resolution STFT loss", ""),
        ("h", "SURVEYED FOR CONTEXT, NOT IMPLEMENTED", "", "", ""),
        ("p", "Wang & Chen, TASLP 2018 (overview)  ·  DCCRN, Interspeech 2020  ·  "
              "FullSubNet, ICASSP 2021  ·  DPRNN, ICASSP 2020  ·  DNS Challenge", "", "", ""),
    ]),
    ("amber", "Datasets", [
        ("e", "Military Audio Dataset", "https://github.com/kaen2891/military_audio_dataset",
         "Kim, Yoon & Jung, Scientific Data 11:668, 2024",
         "CC BY 4.0  ·  primary defence noise  ·  5,655 train / 830 test"),
        ("e", "LibriSpeech", "https://www.openslr.org/12",
         "Panayotov et al., ICASSP 2015", "clean speech for training"),
        ("e", "VCTK-DEMAND", "https://datashare.ed.ac.uk/handle/10283/2791",
         "Valentini-Botinhao et al., Interspeech 2016", "held-out evaluation only  ·  never trained on"),
        ("e", "DEMAND noise database", "https://zenodo.org/records/1227121",
         "Thiemann, Ito & Vincent, 2013", "the noise half of VCTK-DEMAND"),
        ("e", "ESC-50", "https://github.com/karolpiczak/ESC-50",
         "Piczak, ACM Multimedia 2015", "CC BY-NC  ·  siren and wind  ·  evaluation only"),
        ("h", "LICENSING", "", "", ""),
        ("p", "Every dataset licence verified. ESC-50 is CC BY-NC, so it is used for evaluation "
              "only and excluded from every deployment-facing claim.", "", "", ""),
    ]),
    ("green", "Metrics, Standards and Baseline", [
        ("e", "PESQ", "https://www.itu.int/rec/T-REC-P.862",
         "ITU-T P.862 / P.862.2  ·  Rix et al., ICASSP 2001", "perceptual quality  ·  target > 2.5"),
        ("e", "STOI",
         sch("An Algorithm for Intelligibility Prediction of Time-Frequency Weighted Noisy Speech Taal"),
         "Taal, Hendriks, Heusdens & Jensen, IEEE TASLP 19(7), 2011", "intelligibility  ·  target > 0.85"),
        ("e", "SI-SNR", "https://arxiv.org/abs/1811.02508",
         "Le Roux, Wisdom, Erdogan & Hershey, ICASSP 2019", "target > 15 dB"),
        ("e", "ITU-T G.114", "https://www.itu.int/rec/T-REC-G.114",
         "One-way transmission time", "the 150 ms budget our 83.6 ms is measured against"),
        ("e", "Spectral subtraction",
         sch("Suppression of acoustic noise in speech using spectral subtraction Boll 1979"),
         "Boll, IEEE TASSP 27(2), 1979", "our classical baseline  ·  +0.57 on engines, +0.07 on gunfire"),
        ("h", "TOOLCHAIN", "", "", ""),
        ("p", "PyTorch 2.13  ·  ONNX Runtime 1.29 (opset 11)  ·  onnx-simplifier 0.7.3  ·  "
              "silero-vad 6.2.1  ·  pesq / pystoi  ·  Raspberry Pi 5", "", "", ""),
    ]),
]

SW, SH = 13.333, 7.5
M, GAP, TOP = 0.42, 0.26, 1.02
CW = (SW - 2 * M - 2 * GAP) / 3

H_ENTRY, H_CITE, H_NOTE, H_HDR = 0.26, 0.20, 0.19, 0.40
GAP_ENTRY = 0.10


def wrap_lines(text, width_in, pt):
    """Rough line count: average glyph is about 0.5 x point size for Calibri."""
    chars = max(10, int(width_in * 72 / (pt * 0.50)))
    n, cur = 1, 0
    for w in text.split():
        add = len(w) + (1 if cur else 0)
        if cur + add > chars:
            n, cur = n + 1, len(w)
        else:
            cur += add
    return n


def card_height(rows, tw):
    y = 0.52
    for kind, a1, _, a3, a4 in rows:
        if kind == "h":
            y += H_HDR
        elif kind == "p":
            y += wrap_lines(a1, tw, 9.5) * 0.165 + 0.10
        else:
            y += wrap_lines(a1, tw, 13) * H_ENTRY
            y += wrap_lines(a3, tw, 10) * H_CITE
            if a4:
                y += wrap_lines(a4, tw, 9) * H_NOTE
            y += GAP_ENTRY
    return y + 0.14


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def style(run, size, colour, bold=False, url=None):
    run.font.size = Pt(size)
    run.font.name = "Calibri"
    run.font.bold = bold
    if url:
        run.hyperlink.address = url
    run.font.color.rgb = colour          # set after the link, so it wins over theme blue
    run.font.underline = False


prs = Presentation()
prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

tf = textbox(slide, M, 0.22, SW - 2 * M, 0.55)
r = tf.paragraphs[0].add_run()
r.text = "RESEARCH AND REFERENCES"
r.font.size, r.font.name, r.font.bold, r.font.color.rgb = Pt(30), "Cambria", True, INK
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

tw = CW - 0.34
CH = max(card_height(rows, tw) for _, _, rows in CARDS)

for i, (theme, head, rows) in enumerate(CARDS):
    t = THEME[theme]
    x = M + i * (CW + GAP)

    body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(TOP), Inches(CW), Inches(CH))
    body.adjustments[0] = 0.035
    body.fill.solid(); body.fill.fore_color.rgb = t["body"]
    body.line.color.rgb = t["line"]; body.line.width = Pt(1)
    body.shadow.inherit = False

    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x + 0.12), Inches(TOP + 0.08), Inches(CW - 0.24), Inches(0.36))
    pill.adjustments[0] = 0.28
    pill.fill.solid(); pill.fill.fore_color.rgb = t["pill"]
    pill.line.color.rgb = t["line"]; pill.line.width = Pt(1)
    pill.shadow.inherit = False
    ptf = pill.text_frame
    ptf.margin_left = ptf.margin_right = ptf.margin_top = ptf.margin_bottom = 0
    ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
    pr = ptf.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
    style(pr.add_run(), 13, t["dark"], bold=True)
    pr.runs[0].text = head

    y, tx = TOP + 0.52, x + 0.17
    for kind, a1, a2, a3, a4 in rows:
        if kind == "h":
            ln = slide.shapes.add_connector(1, Inches(tx), Inches(y + 0.04),
                                            Inches(tx + tw), Inches(y + 0.04))
            ln.line.color.rgb = t["line"]; ln.line.width = Pt(0.75)
            f = textbox(slide, tx, y + 0.11, tw, 0.22)
            style(f.paragraphs[0].add_run(), 10, t["dark"], bold=True)
            f.paragraphs[0].runs[0].text = a1
            y += H_HDR
        elif kind == "p":
            n = wrap_lines(a1, tw, 9.5)
            f = textbox(slide, tx, y, tw, n * 0.165 + 0.06)
            style(f.paragraphs[0].add_run(), 9.5, MUTED)
            f.paragraphs[0].runs[0].text = a1
            y += n * 0.165 + 0.10
        else:
            n1 = wrap_lines(a1, tw, 13)
            n3 = wrap_lines(a3, tw, 10)
            n4 = wrap_lines(a4, tw, 9) if a4 else 0
            h = n1 * H_ENTRY + n3 * H_CITE + n4 * H_NOTE
            f = textbox(slide, tx, y, tw, h + 0.06)
            p0 = f.paragraphs[0]
            style(p0.add_run(), 13, t["dark"], bold=True, url=a2)
            p0.runs[0].text = a1
            p1 = f.add_paragraph()
            style(p1.add_run(), 10, MUTED)
            p1.runs[0].text = a3
            if a4:
                p2 = f.add_paragraph()
                style(p2.add_run(), 9, FAINT)
                p2.runs[0].text = a4
            y += h + GAP_ENTRY

slide.notes_slide.notes_text_frame.text = (
    "Every paper, dataset and standard name is a live hyperlink. Papers whose DOI could not "
    "be verified link to a Google Scholar search rather than a guessed DOI."
)

out = Path("docs/references_slide.pptx")
prs.save(out)
print(f"wrote {out} ({out.stat().st_size} bytes), card height {CH:.2f}in of {SH}in slide")
